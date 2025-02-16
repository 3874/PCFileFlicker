import os, re, json, threading, qrcode, socket
import shutil, requests, webbrowser, sys, xlrd, openpyxl
import tkinter as tk
from tkinter import filedialog, messagebox, ttk, simpledialog, scrolledtext
from docx import Document
from PyPDF2 import PdfReader
from tinydb import TinyDB, Query
from bs4 import BeautifulSoup
from qrserver import qr_app
from pptx import Presentation

if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

QR_settings = None
dashboard_window = None
OPENAI_API = None
config_path = None
QR_PORT = 9000
QR_VERSION = '0.0.5'

UPLOAD_FOLDER = os.path.join(BASE_DIR, 'drive')
DB_FOLDER = os.path.join(BASE_DIR, 'db')
SETTING_FOLDER = os.path.join(BASE_DIR, 'setting')

for folder in [UPLOAD_FOLDER, DB_FOLDER, SETTING_FOLDER]:
    if not os.path.exists(folder):
        os.makedirs(folder)

def save_settings(settings):
    
    with open(config_path, "w") as json_file:
        json.dump(settings, json_file, indent=4)
        
    messagebox.showinfo("저장 완료", "설정이 JSON 형태로 저장되었습니다.")

config_path = os.path.join(SETTING_FOLDER, "config.json")
if not os.path.exists(config_path):
    settings = {
        "port": QR_PORT,
        "OPENAI_KEY": ''
    }
    save_settings(settings)
    
def load_settings():
    
    if os.path.exists(config_path):
        with open(config_path, "r") as json_file:
            settings = json.load(json_file)
            return settings
    else:
        return None

QR_settings = load_settings()

def open_api_settings(OPENAI_API_VALUE=None):
    user_input = simpledialog.askstring("CONFIG","Input OPENAI api key.", initialvalue=OPENAI_API_VALUE)

    if user_input is not None:
        OPENAI_API = user_input
        settings = {
            "port": QR_PORT,
            "OPENAI_KEY": OPENAI_API
        }
        save_settings(settings)
    else:
        messagebox.showinfo("입력이 취소되었습니다.")

db_path = os.path.join(DB_FOLDER,'fileDB.json')

db = TinyDB(db_path)

fileTable = db.table('fileTable')

ITEMS_PER_PAGE = 40
current_page = 1
   
def get_local_ip():
    hostname = socket.gethostname()
    local_ip = socket.gethostbyname(hostname)
    return local_ip

def generate_qr_code(data):
    img = qrcode.make(data)
    img.show()
    
def open_server_mgt_dashboard():
    global dashboard_window

    if dashboard_window is not None and dashboard_window.winfo_exists():
        dashboard_window.focus_force()
        return
    
    dashboard_window = tk.Toplevel(root)
    dashboard_window.title("Server Manager")
    dashboard_window.geometry("300x150")

    # Title Label
    title_label = tk.Label(dashboard_window, text="Server Manager", font=("Arial", 16, "bold"))
    title_label.pack(pady=10) 

    button_frame = tk.Frame(dashboard_window)
    button_frame.pack(side=tk.BOTTOM, pady=10) 
    
    server_state = False 
    
    private_ip = get_local_ip()
    port= QR_settings['port']
    
    def show_qrcode(targetURL):
            
        if not server_state:
            messagebox.showwarning("Warning", "Server is OFF, please start the server to generate the QR code.")
            dashboard_window.focus_force() 
            return
    
        if targetURL:
            generate_qr_code(targetURL)
        else:
            messagebox.showwarning("Warning", "Please enter a valid IP address.")

    def run_server():
        qr_app.run(host='0.0.0.0', port=port, use_reloader=False)

    def toggle_server():
        nonlocal server_state
        server_state = not server_state 
        private_url = None 

        if server_state:
            server_button.config(text='Turn Off Server', bg='lightcoral')

            try:
                server_thread = threading.Thread(target=run_server, daemon=True)  # Run server in a daemon thread
                server_thread.start()

                private_url = f"http://{private_ip}:{port}/"
                messagebox.showinfo("Server Status", f"Server is now ON. Access it at http://{private_ip}:{port}/")
                show_qrcode(private_url)
                
            except Exception as e:
                messagebox.showerror("Error", f"Failed to start the server: {str(e)}")
                server_state = False
                server_button.config(text='Turn On Server', bg='lightgreen')
        else:
            server_button.config(text='Turn On Server', bg='lightgreen')

        dashboard_window.focus_force() 

    server_button = tk.Button(button_frame, text='Turn On Server', command=toggle_server, bg='lightgreen')
    server_button.pack(side=tk.LEFT, padx=5)
       
def get_latest_version():
    url = 'https://s3.us-west-1.amazonaws.com/fileflicker.me/download/current_version.txt'
    response = requests.get(url)
    return response.text.strip()

def check_for_update():
    current_version = QR_VERSION
    latest_version = get_latest_version()
    
    compare_result = compare_versions(current_version, latest_version)
    
    if compare_result:
        update_message = f"New version available: {latest_version}\nDo you want to update?"
        if messagebox.askyesno("Update Available", update_message):
            new_file = download_latest_version()
            messagebox.showinfo("Update", f"{new_file}: Download completed.")
        else:
            messagebox.showinfo("Update", "You are using the current version.")
    else:
        messagebox.showinfo("No Update", "You are using the latest version.")

def compare_versions(current_version, new_version):

    current_parts = list(map(int, current_version.split('.')))
    new_parts = list(map(int, new_version.split('.')))

    for current, new in zip(current_parts, new_parts):
        if current < new:
            return True
        elif current > new:
            return False
    
    return False

def download_latest_version():
    url = "https://s3.us-west-1.amazonaws.com/fileflicker.me/download/latest_update.zip"
    file_name = "latest_update.zip"
    output_path = os.path.join(os.getcwd(), file_name)

    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()  # 요청이 성공적인지 확인 (HTTP 오류 체크)
        
        with open(output_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        
        return output_path  # 다운로드 완료 후 파일 경로 반환

    except requests.exceptions.HTTPError as http_err:
        messagebox.showerror("Download Error", f"HTTP error occurred: {http_err}")
    except Exception as err:
        messagebox.showerror("Download Error", f"An error occurred: {err}")

    return None 

def google_search_scrape(query, num_results=30):
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
    }
    search_url = f"https://www.google.com/search?q={query}&num={num_results}"
    
    response = requests.get(search_url, headers=headers)
    soup = BeautifulSoup(response.text, 'html.parser')
    
    search_results = []
    for g in soup.find_all('div', class_='g'):
        title = g.find('h3').text if g.find('h3') else 'No title'
        a_tag = g.find('a')
        link = a_tag['href'] if a_tag and 'href' in a_tag.attrs else 'No link'
        snippet = g.find('span', class_='aCOpRe').text if g.find('span', class_='aCOpRe') else 'No snippet'
        
        search_results.append({
            'title': title,
            'url': link,
            'snippet': snippet
        })
        
    return search_results
    
def get_prompt(promptNo):
    
    prompt_path = os.path.join(SETTING_FOLDER, "prompt.json")
    with open(prompt_path, 'r') as file:
        variables = json.load(file)

    prompt = variables[promptNo]
    return prompt

def delete_file(file_id, file_path):
    if messagebox.askyesno("Confirm Deletion", "Are you sure you want to delete this file?"):
        try:
            fileTable.remove(doc_ids=[file_id])
            os.remove(file_path)
            messagebox.showinfo("Success", "File deleted successfully")
            refresh_file_list()
            return True
        except Exception as e:
            messagebox.showerror("Error", f"Failed to delete file: {str(e)}")
            return False
    return False


def update_file_info(file_id, company, company_code, tags, summary, comments):
    fileTable.update({'company': company, 'company_code': company_code, 'tags': tags, 'summary': summary, 'comments': comments}, doc_ids=[file_id])

def update_DB_with_bulk_file():

    file_list = os.listdir(UPLOAD_FOLDER)
    file_list = [f for f in file_list if os.path.isfile(os.path.join(UPLOAD_FOLDER, f))]
    
    if not file_list:
        messagebox.showinfo("No Files", "There are no files in the upload folder.")
        return
    
    with TinyDB(db_path) as db:
        fileTable = db.table('fileTable')

        db_files = fileTable.all()
        db_file_list = [item['file_name'] for item in db_files]
        save_path = None
        
        files_updated = False

        for file_name in file_list:
            if file_name not in db_file_list:
                save_path = os.path.join(UPLOAD_FOLDER, file_name)
                fileTable.insert({'file_name': file_name.encode('utf-8').decode('utf-8'), 'summary': '', 'location': save_path, 'company': '', 'company_code': '', 'tags': '', 'comments': '' })
                messagebox.showinfo("Success", f"{file_name} is updated.")
                files_updated = True

        if not files_updated:
            messagebox.showinfo("No Updates", "No new files were added to the database.")

    db.close()
    
def Json2String(json_obj, indent=0):
    result = ""
    indent_str = " " * (indent * 4)
    companyName = None
    tags = []

    if isinstance(json_obj, dict):
        for key, value in json_obj.items():
            if isinstance(value, dict):
                nested_result, nested_company_name, nested_tags = Json2String(value, indent + 1)
                result += f"{indent_str}{key}:\n{nested_result}"
                if nested_company_name:
                    companyName = nested_company_name
                if nested_tags:
                    tags = nested_tags
            elif isinstance(value, list):
                if key == "tags" and isinstance(value, list):
                    tags = value
                result += f"{indent_str}{key}: {value}\n"
            else:
                if key == "company_name":
                    companyName = value
                result += f"{indent_str}{key}: {value}\n"
    else:
        result += f"{indent_str}{json_obj}\n"

    return result, companyName, tags

def process_and_store_file(file_path):
    try:
        file_name = os.path.basename(file_path)
        save_path = os.path.join(UPLOAD_FOLDER, file_name)
        File = Query()
        if fileTable.contains(File.file_name == file_name):
            messagebox.showwarning("Duplicate File", f"The file '{file_name}' already exists in the database.")
            return None, None

        text = extract_text_from_file(file_path)
        prompt =  get_prompt('automatic_company_analysis')
        summary = ai_analysis(prompt, text)
        shutil.copy(file_path, save_path)

        fileTable.insert({
            'file_name': file_name,
            'summary': summary,
            'location': save_path,
            'company': '',
            'company_code': '',
            'tags': '',
            'comments': ''
        })
        
        return file_name, summary
    except Exception as e:
        messagebox.showerror("Error", str(e))
        return None, None

def extract_json_from_string(text):
    json_pattern = re.compile(r'\{(?:[^{}]|(?R))*\}')
    matches = json_pattern.findall(text)
    json_objects = []
    
    for match in matches:
        try:
            json_obj = json.loads(match)
            json_objects.append(json_obj)
        except json.JSONDecodeError:
            continue
    
    return json_objects

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ''
    
    if ext == '.pdf':
        # PDF 파일에서 텍스트 추출
        reader = PdfReader(file_path)
        for page in reader.pages:
            text += page.extract_text() + '\n'
    
    elif ext in ['.doc', '.docx']:
        # Word 파일에서 텍스트 추출
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            text += paragraph.text + '\n'
    
    elif ext == '.txt':
        # 텍스트 파일에서 텍스트 읽기
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    
    elif ext == '.pptx':
        # PPTX 파일에서 텍스트 추출
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
    
    elif ext == '.ppt':
        # PPT 파일에서 텍스트 추출할 수 없음을 경고
        raise ValueError("PPT format is not supported. Please convert to PPTX format.")
    
    elif ext in ['.xls', '.xlsx']:
        # XLS/XLSX 파일에서 텍스트 추출
        if ext == '.xls':
            workbook = xlrd.open_workbook(file_path)
            for sheet in workbook.sheets():
                for row_idx in range(sheet.nrows):
                    row = sheet.row(row_idx)
                    text += ' '.join([str(cell.value) for cell in row]) + '\n'
        else:  # .xlsx
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
    
    return text

def ai_analysis(prompt, text):

    url = 'https://api.openai.com/v1/chat/completions'
    
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {OPENAI_API}'
    }

    data = {
        'model': "gpt-4o-mini",
        'messages': [
            {"role": "system", "content": "You are a helpful assistant that summarizes text."},
            {"role": "user", "content": f"{prompt}\n\n{text}"}
        ],
        'max_tokens': 4096,
        'temperature': 0.7
    }

    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()
        result = response.json()
        summary = result['choices'][0]['message']['content'].strip()
        return summary
    except requests.exceptions.HTTPError as err:
        error_details = response.json()


def select_and_process_file():
    file_path = filedialog.askopenfilename(
        filetypes=[("All Files", "*.*"), ("PDF Files", "*.pdf"), ("Word Documents", "*.doc;*.docx"), ("Text Files", "*.txt")]
    )
    if file_path:
        def process():
            file_name, summary = process_and_store_file(file_path)
            if file_name and summary:
                root.after(0, lambda: messagebox.showinfo("Success", f"File processed and stored successfully.\n\nFile Name: {file_name}\nSummary: {summary}"))
                root.after(0, refresh_file_list)
        
        run_with_spinner(root, process)
        
def append_data(json_window, summary_text, json_str):
    summary_text.insert(tk.END, "\n\n +++++ New Results +++++\n\n")
    summary_text.insert(tk.END, json_str + "\n")
    json_window.destroy()

def extract_text_from_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        # 응답의 내용을 파싱
        soup = BeautifulSoup(response.text, 'html.parser')

        # CSS, JavaScript, HTML 태그 제거
        for script_or_style in soup(['script', 'style']):
            script_or_style.decompose()

        # 페이지의 모든 텍스트 내용을 추출
        page_text = soup.get_text(separator='\n', strip=True)

        return page_text

    except requests.exceptions.RequestException as e:
        messagebox.ERROR("Error", f"요청 중 오류가 발생했습니다: {e}")
        return None

def show_json_window(data, summary_text, append_btn=True, analyze_btn=True):
    json_window = tk.Toplevel(root)
    json_window.title("Searching Results")
    json_window.geometry("800x600")

    # 텍스트 위젯을 담을 프레임 생성
    text_frame = tk.Frame(json_window)
    text_frame.pack(expand=True, fill=tk.BOTH, padx=10, pady=10)

    # 스크롤바 생성
    scrollbar = tk.Scrollbar(text_frame)
    scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

    # 텍스트 위젯 생성
    json_text = tk.Text(text_frame, wrap=tk.WORD, yscrollcommand=scrollbar.set)
    json_text.pack(expand=True, fill=tk.BOTH)

    scrollbar.config(command=json_text.yview)

    try:
        # JSON 데이터 파싱 시도
        if isinstance(data, str):
            parsed_data = json.loads(data)
        else:
            parsed_data = data

        if isinstance(parsed_data, list):
            # JSON 배열 처리
            formatted_data = []
            for item in parsed_data:
                if isinstance(item, dict):
                    formatted_data.append(json.dumps(item, indent=4, ensure_ascii=False).encode('utf-8').decode())
                else:
                    formatted_data.append(str(item))
            json_str = "\n\n".join(formatted_data)
        elif isinstance(parsed_data, dict):
            # 단일 JSON 객체 처리
            json_str = json.dumps(parsed_data, indent=4, ensure_ascii=False).encode('utf-8').decode()
        else:
            # 기타 데이터 타입
            json_str = str(parsed_data)
    except json.JSONDecodeError:
        # JSON 파싱 실패 시 원본 데이터 표시
        json_str = str(data)

    json_text.insert(tk.END, json_str)
    
    def close_window():
        json_window.destroy()
        
    def smart_analyze():
        pattern = r'\{[^}]+\}'
        matches = re.findall(pattern, json_str)
        result = None

        for match in matches:
            try:
                item_dict = json.loads(match)
                if 'url' in item_dict:
                    ExtractDict = extract_text_from_url(item_dict['url'])
                    prompt = f"First, summarize it based on {ExtractDict}."
                    text = f"Then, compare this summary with {summary_text} to determine if this is relevant or not. if it is relevant, response summary; if not, response with 'no relevent' and 'url'."
                    if result:
                        result += ai_analysis(prompt, text)
                    else: 
                        result = ai_analysis(prompt, text)
                        
            except json.JSONDecodeError as e:
                messagebox.ERROR("error",f"Error decoding JSON: {e}")

        json_text.insert(tk.END, "\n\n +++++ New Results +++++\n\n")
        json_text.insert(tk.END, result + "\n")

    button_frame = tk.Frame(json_window)
    button_frame.pack(pady=10)

    if append_btn:
        append_button = tk.Button(button_frame, text="Append", command=lambda: append_data(json_window, summary_text, json_str), bg="red", fg="white")
        append_button.pack(side=tk.LEFT, padx=5)

    if analyze_btn:
        analyze_button = tk.Button(button_frame, text="Smart analyze", command=smart_analyze, bg="black", fg="white")
        analyze_button.pack(side=tk.LEFT, padx=5)

    close_button = tk.Button(button_frame, text="Close", command=close_window)
    close_button.pack(side=tk.LEFT, padx=5)
    
    # 모달 설정
    json_window.transient(root)
    json_window.grab_set()
    root.wait_window(json_window)
    
def get_file_list():
    filedata = fileTable.all()
    data = list(reversed(filedata))
    return data

def refresh_file_list(page=1):
    global current_page
    current_page = page
    for i in tree.get_children():
        tree.delete(i)
    files = get_file_list()
    
    start_index = (page - 1) * ITEMS_PER_PAGE
    end_index = start_index + ITEMS_PER_PAGE
    paged_files = files[start_index:end_index]

    for file in paged_files:
        tree.insert('', 'end', values=(file.doc_id, file['file_name'], file.get('tags', ''), file.get('summary', '')))
    
    update_pagination_controls('file')

def update_pagination_controls(classify):
    # Clear existing buttons
    for btn in pagination_frame.winfo_children():
        btn.destroy()

    if classify == 'file': 
        total_files = len(get_file_list())
    else:
        total_files = 0
    
    total_pages = (total_files + ITEMS_PER_PAGE - 1) // ITEMS_PER_PAGE

    prev_btn = None  
    next_btn = None  
    first_btn = None 
    last_btn = None  

    # First Button
    if total_pages > 0:
        first_btn = tk.Button(pagination_frame, text='≪', command=lambda: refresh_file_list(1))
        first_btn.pack(side='left')
    
    # Previous Button
    if current_page > 1:
        prev_btn = tk.Button(pagination_frame, text='＜', command=lambda: refresh_file_list(current_page - 1))
        prev_btn.pack(side='left')

    # Current Page Label
    page_label = tk.Label(pagination_frame, text=f'{current_page} of {total_pages}')
    page_label.pack(side='left')

    # Next Button
    if current_page < total_pages:
        next_btn = tk.Button(pagination_frame, text='＞', command=lambda: refresh_file_list(current_page + 1))
        next_btn.pack(side='left')

    # Last Button
    if total_pages > 0:
        last_btn = tk.Button(pagination_frame, text='≫', command=lambda: refresh_file_list(total_pages))
        last_btn.pack(side='left')

    # Disable/Enable buttons based on the current page
    if prev_btn is not None:
        prev_btn.config(state=tk.NORMAL if current_page > 1 else tk.DISABLED)
    
    if next_btn is not None:
        next_btn.config(state=tk.NORMAL if current_page < total_pages else tk.DISABLED)

    if first_btn is not None:
        first_btn.config(state=tk.NORMAL if current_page > 1 else tk.DISABLED)

    if last_btn is not None:
        last_btn.config(state=tk.NORMAL if current_page < total_pages else tk.DISABLED)
        
def open_file(file_path):
    try:
        if os.path.exists(file_path):
            if os.name == 'nt':  
                os.startfile(file_path)
            elif os.name == 'posix': 
                os.system(f'open "{file_path}"')
            else:
                raise OSError("Unsupported operating system")
        else:
            raise FileNotFoundError(f"The file does not exist: {file_path}")
    except Exception as e:

        root = tk.Tk()
        root.withdraw()  
        messagebox.showerror("Error", f"Failed to open file: {str(e)}")

def get_prompt_from_user(parent_window, prompts, default_index=0):

    def on_select():
        nonlocal selected_prompt
        selected_prompt = combobox.get()
        dialog.destroy()

    selected_prompt = None
    dialog = tk.Toplevel(parent_window)
    dialog.title("Select Prompt")
    dialog.geometry("300x200")
    
    def update_combobox_width():
        # Get the current width of the dialog window
        dialog_width = dialog.winfo_width()
        # Adjust combobox width (subtracting padding)
        combobox_width = dialog_width - 20  # 10px padding on each side
        combobox.config(width=combobox_width)
    
    tk.Label(dialog, text="Select a prompt:").pack(padx=10, pady=5)
    combobox = ttk.Combobox(dialog, values=prompts)
    combobox.pack(padx=10, pady=5)
    
    if prompts:
        combobox.current(default_index)
    
    ttk.Button(dialog, text="OK", command=on_select).pack(pady=10)
    dialog.bind("<Configure>", lambda e: update_combobox_width())
    dialog.wait_window()

    return selected_prompt

def open_url(url):
    webbrowser.open(url) 
    
def show_file_details(event):
    item = tree.selection()[0]
    file_id = int(tree.item(item, "values")[0])
    file_details = fileTable.get(doc_id=file_id)

    if file_details:
        detail_window = tk.Toplevel(root)
        detail_window.title("File Details")
        detail_window.geometry("800x600")

        file_name_frame = tk.Frame(detail_window)
        file_name_frame.pack(fill='x', padx=10, pady=5)
        tk.Label(file_name_frame, text="File Name:", font=('Arial', 12, 'bold')).pack(side='left')
        tk.Label(file_name_frame, text=file_details['file_name'], wraplength=380).pack(side='left', padx=5)
        
        def delete_and_close():
            if delete_file(file_id, file_details['location']):
                detail_window.destroy()

        delete_button = tk.Button(file_name_frame, text="Delete", command=delete_and_close, bg="red", fg="white")
        delete_button.pack(side='left', padx=5)

        tk.Label(detail_window, text="File Location:", font=('Arial', 12, 'bold')).pack(anchor='w', padx=10, pady=5)
        location_label = tk.Label(detail_window, text=file_details['location'], wraplength=480, cursor="hand2", fg="blue")
        location_label.pack(anchor='w', padx=10)
        location_label.bind("<Button-1>", lambda e: open_file(file_details['location']))

        tags_frame = tk.Frame(detail_window)
        tags_frame.pack(fill='x', padx=10, pady=5)
        tk.Label(tags_frame, text="Tags:", font=('Arial', 12, 'bold')).pack(side='left')
        tags_entry = tk.Entry(tags_frame, width=30)
        tags_entry.pack(fill=tk.BOTH, padx=(5, 0))
        tags_entry.insert(0, file_details.get('tags', ''))

        def analyze_file():
            
            prompts = [
                "summarize this in primary language:",
                "summarize this in Korean:",
                "extract all the tags from the given file and create its comma-separated list in primary language:",
                "Manual input"
            ]
                
            prompt = get_prompt_from_user(detail_window, prompts, default_index=0)
            
            if not prompt:
                detail_window.focus()
                messagebox.showinfo("Info", "Analyzing canceled. No prompt provided.")
                return 
            
            def AIanalyze():
                try:
                    text = extract_text_from_file(file_details['location'])
                    new_summary = ai_analysis(prompt, text)
                except Exception as e:
                    messagebox.showerror("Error", f"An error occurred: {e}")
                    return

                def update_ui():
                    summary_text.insert(tk.END, "\n\n--- New Results ---\n\n")
                    summary_text.insert(tk.END, new_summary)
                    summary_text.see(tk.END)
                
                detail_window.after(0, update_ui)
                    
            run_with_spinner(detail_window, AIanalyze)

        summary_frame = tk.Frame(detail_window)
        summary_frame.pack(fill='x', padx=10, pady=5)
        
        tk.Label(summary_frame, text="Summary:", font=('Arial', 12, 'bold')).pack(side='left')

        aianalyze_button = tk.Button(summary_frame, text="Search in file", command=analyze_file, bg="orange", fg="white")
        aianalyze_button.pack(side='left', padx=10)
        
        summary_text = tk.Text(detail_window, wrap=tk.WORD, height=10)
        summary_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        summary_text.insert(tk.END, file_details.get('summary', ''))

        tk.Label(detail_window, text="Comments:", font=('Arial', 12, 'bold')).pack(anchor='w', padx=10, pady=5)
        comments_text = tk.Text(detail_window, wrap=tk.WORD, height=3)
        comments_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        comments_text.insert(tk.END, file_details.get('comments', ''))

        def save_details():
            company = ''
            company_code = ''
            tags = tags_entry.get()
            summary = summary_text.get("1.0", tk.END).strip()
            comments = comments_text.get("1.0", tk.END).strip()
            update_file_info(file_id, company, company_code, tags, summary, comments)
            messagebox.showinfo("Success", "File details updated successfully")
            detail_window.destroy()
            refresh_file_list()
        
        def close_window():
            detail_window.destroy()
        
        def open_mail_client(file_details):
            subject = "File Details: " + file_details['file_name']
            body = f"File Name: {file_details['file_name']}\nLocation: {file_details['location']}\nCompany: {file_details.get('company', '')}\nCode: {file_details.get('company_code', '')}\nTags: {file_details.get('tags', '')}\nSummary: {file_details.get('summary', '')}\nComments: {file_details.get('comments', '')}"
            body = body.replace("\n", "%0D%0A")
            mailto_link = f"mailto:?subject={subject}&body={body}"
            webbrowser.open(mailto_link)
            messagebox.showinfo("Attach File", "Please attach the file manually in your email client.")

        button_frame = tk.Frame(detail_window)
        button_frame.pack(pady=10)

        save_button = tk.Button(button_frame, text="Save", command=save_details, bg="blue", fg="white")
        save_button.pack(side=tk.LEFT, padx=5)

        close_button = tk.Button(button_frame, text="Close", command=close_window)
        close_button.pack(side=tk.LEFT, padx=5)

        mailto_button = tk.Button(button_frame, text="Email", command=lambda: open_mail_client(file_details), bg="green", fg="white")
        mailto_button.pack(side=tk.LEFT, padx=5)
        
        selected_option = tk.StringVar(root)
        selected_option.set("옵션 1")

        detail_window.transient(root)
        detail_window.grab_set()
        root.wait_window(detail_window)

def search_files():
    search_term = search_entry.get().strip().lower()
    if not search_term:
        refresh_file_list()
        return

    for i in tree.get_children():
        tree.delete(i)

    def insert_into_tree(results, columns):
        for file in results:
            values = [file.doc_id] + [file.get(col, '') for col in columns]
            tree.insert('', 'end', values=values)

    def test_function(value):
        return value is not None and search_term in value.lower()

    File = Query()
    results = fileTable.search(
        (File.file_name.test(test_function)) |
        (File.summary.test(test_function)) |
        (File.tags.test(test_function))
    )
    columns = ['file_name', 'tags', 'summary']
    insert_into_tree(results, columns)
    
    search_entry.delete(0, 'end')

def create_spinner_window(parent):
    spinner_window = tk.Toplevel(parent)
    spinner_window.overrideredirect(True)
    spinner_window.attributes('-topmost', True)
    spinner_window.geometry('200x100')
    
    parent_x = parent.winfo_x()
    parent_y = parent.winfo_y()
    parent_width = parent.winfo_width()
    parent_height = parent.winfo_height()
    
    x = parent_x + (parent_width // 2) - 100
    y = parent_y + (parent_height // 2) - 50
    
    spinner_window.geometry(f'+{x}+{y}')
    
    spinner_frame = tk.Frame(spinner_window, bg='white')
    spinner_frame.place(relx=0.5, rely=0.5, anchor='center')
    
    spinner = ttk.Progressbar(spinner_frame, mode='indeterminate', length=100)
    spinner.pack(pady=10)
    
    label = tk.Label(spinner_frame, text="Processing...", bg='white')
    label.pack()
    
    spinner.start(10)
    return spinner_window

def start_spinner(parent):
    return create_spinner_window(parent)

def stop_spinner(spinner_window):
    if spinner_window:
        spinner_window.destroy()

def run_with_spinner(parent, func, *args):
    spinner_window = start_spinner(parent)
    
    def wrapper():
        try:
            func(*args)
        finally:
            parent.after(0, lambda: stop_spinner(spinner_window))
    
    threading.Thread(target=wrapper).start()
    
def open_webpage(url):
    webbrowser.open(url)
    
def create_submenu(menu_bar):
    file_menu = tk.Menu(menu_bar, tearoff=0)
    menu_bar.add_cascade(label="FileFlicker", menu=file_menu)
    file_menu.add_command(label="About", command=lambda: open_webpage('https://francishan.notion.site/About-QR-File-Manager-cd8b41d3c31c454f91db95df842d792d'))
    
    upload_menu = tk.Menu(menu_bar, tearoff=0)
    menu_bar.add_cascade(label="Add", menu=upload_menu)
    upload_menu.add_command(label="Single file", command=select_and_process_file)
    upload_menu.add_command(label="Unregistered file", command=update_DB_with_bulk_file)
    
    server_menu = tk.Menu(menu_bar, tearoff=0)
    menu_bar.add_cascade(label="Web Access", menu=server_menu)
    server_menu.add_command(label="Server Manager", command=open_server_mgt_dashboard)
    server_menu.add_command(label="ReadMe", command=lambda: open_webpage('https://francishan.notion.site/On-my-mobile-bacf8d67c55b4543a5cd4c2313b94371'))

    help_menu = tk.Menu(menu_bar, tearoff=0)
    menu_bar.add_cascade(label="Help", menu=help_menu)
    help_menu.add_command(label="How to use", command=lambda: open_webpage('https://francishan.notion.site/How-to-Use-c8e131ec3a35417786d3f03556e6f741'))
    help_menu.add_command(label="Report Issue", command=lambda: open_webpage('https://docs.google.com/forms/d/13aoyEZhTE3N9M82gWQxRo9Ir8_s3s7YXeQq8rF3-E88'))
    
    help_menu.add_separator()    
    new_menu = tk.Menu(file_menu, tearoff=0)
    help_menu.add_cascade(label="Settings", command=lambda: open_api_settings(OPENAI_API))

    help_menu.add_command(label="Check for Updates", command=check_for_update)
        
def show_readme_window(contents):
    readme_window = tk.Toplevel(root)
    readme_window.title("Readme")
    readme_window.geometry("600x400")

    text_area = scrolledtext.ScrolledText(readme_window, wrap=tk.WORD)
    text_area.insert(tk.INSERT, contents)
    text_area.configure(state='disabled')
    text_area.pack(expand=True, fill='both')
        
def close_databases():
    db.close()
    messagebox.showinfo("Close","Databases closed.")
        
if __name__ == "__main__":
    try:
        root = tk.Tk()
        root.title("FileFlicker")
        root.state('zoomed')

        screen_ratio = 0.8
        screen_width = int(root.winfo_screenwidth()*screen_ratio)

        menu_bar = tk.Menu(root)
        root.config(menu=menu_bar)
        create_submenu(menu_bar)

        search_frame = tk.Frame(root)
        search_frame.pack(pady=10)

        search_entry = tk.Entry(search_frame, width=50)
        search_entry.pack(side=tk.LEFT, padx=5)

        search_button = tk.Button(search_frame, text="▶", command=search_files)
        search_button.pack(side=tk.LEFT)

        search_entry.bind("<Return>", lambda event: search_files())

        tree = ttk.Treeview(root, columns=('ID', 'File Name', 'Tags', 'Summary'), show='headings')
        tree.heading('ID', text='ID')
        tree.heading('File Name', text='File Name')
        tree.heading('Tags', text='Tags')
        tree.heading('Summary', text='Summary')
        tree.column('ID', width=int(screen_width*0.05))
        tree.column('File Name', width=int(screen_width*0.2))
        tree.column('Tags', width=int(screen_width*0.25))
        tree.column('Summary', width=int(screen_width*0.35))
        tree.pack(expand=True, fill='both')
        
        tree.bind("<Double-1>", show_file_details)
        
        select_file_button = tk.Button(root, text="📤 Upload File", command=select_and_process_file)
        select_file_button.pack(side=tk.RIGHT, padx=10)
        
        pagination_frame = tk.Frame(root)
        pagination_frame.pack(side='bottom', pady=10)

        refresh_file_list()
        
        QR_settings = load_settings()

        if not QR_settings or not QR_settings['OPENAI_KEY']:
            messagebox.showwarning("Settings", "No Settings file. put your keys for AI analysis.")
            open_api_settings('')
            QR_settings = load_settings()

        OPENAI_API = QR_settings['OPENAI_KEY']

        root.mainloop()
    except Exception as e:
        messagebox.showerror(f"An error occurred: {e}")
        
    finally:
        close_databases()

