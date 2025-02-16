import logging, os, sys, json, xlrd, openpyxl, requests
#import openai
from flask import Flask, request, jsonify, send_from_directory, render_template
from tinydb import TinyDB
from flask_cors import CORS
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from fuzzywuzzy import fuzz

QRapp = Flask(__name__)
CORS(QRapp)

AI_PROMPT = 'Please summarize the following text in the primary language used:'

if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

TEMPLATE_FOLDER = os.path.join(BASE_DIR, 'templates')
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'drive')
DB_FOLDER = os.path.join(BASE_DIR, 'db')
SETTING_FOLDER = os.path.join(BASE_DIR, 'setting')

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DB_FOLDER, exist_ok=True)
os.makedirs(SETTING_FOLDER, exist_ok=True)

db_path = os.path.join(DB_FOLDER, 'fileDB.json')
db_path2 = os.path.join(DB_FOLDER, 'companyDB.json')

db = TinyDB(db_path)
db2 = TinyDB(db_path2)

fileTable = db.table('fileTable')
compTable = db2.table('compTable')

config_path = os.path.join(SETTING_FOLDER, "config.json")

# config.json 파일이 있는지 확인하고 로드합니다.
if not os.path.exists(config_path):
    raise FileNotFoundError("Config file not found. Ensure config.json exists in the setting directory.")

with open(config_path, "r") as json_file:
    QR_config = json.load(json_file)

if not QR_config or not QR_config.get('OPENAI_KEY'):
    raise ValueError('OpenAI Key is missing in the config file.')

# openai.api_key = QR_config['OPENAI_KEY']

@QRapp.route('/')
def home():
    return render_template('index.html')

@QRapp.route('/gotofile')
def gotofile():
    return render_template('file.html')

@QRapp.route('/gotocompany')
def gotocompany():
    return render_template('company.html')

@QRapp.route('/edit/file/<int:doc_id>', methods=['PUT'])
def edit_file(doc_id):
    data = request.json
    file_entry = fileTable.get(doc_id=doc_id)

    if not file_entry:
        return jsonify({"error": "File not found"}), 404

    # 필드 업데이트
    for key in ['file_name', 'summary', 'tags', 'comments', 'company', 'company_code']:
        if key in data:
            file_entry[key] = data[key]

    # 데이터베이스 업데이트
    fileTable.update(file_entry, doc_ids=[doc_id])

    return jsonify({"success": "File updated", "file": file_entry}), 200

@QRapp.route('/edit/company/<int:corpID>', methods=['PUT'])
def edit_company(corpID):
    data = request.json
    company_entry = compTable.get(doc_id=corpID)

    if not company_entry:
        return jsonify({"error": "Company not found"}), 404

    # 필드 업데이트
    for key in ['company', 'summary', 'tags']:
        if key in data:
            company_entry[key] = data[key]

    # 데이터베이스 업데이트
    compTable.update(company_entry, doc_ids=[corpID])

    return jsonify({"success": "Company updated", "company": company_entry}), 200

@QRapp.route('/files', methods=['GET'])
def get_files():
    start = int(request.args.get('start', 0))
    length = int(request.args.get('length', 10))
    search_value = request.args.get('search[value]', '').lower()

    files = fileTable.all()

    if search_value:
        files = [file for file in files if (
            search_value in file['file_name'].lower() or
            search_value in file.get('summary', '').lower() or
            any(search_value in tag.lower() for tag in file.get('tags', []))
        )]

    records_with_id = [{'doc_id': file.doc_id, **file} for file in files]
    records_with_id.reverse()
    paginated_files = records_with_id[start:start + length]

    response = {
        "draw": int(request.args.get('draw', 1)),
        "recordsTotal": len(fileTable.all()),
        "recordsFiltered": len(records_with_id),
        "data": paginated_files
    }

    return jsonify(response)

@QRapp.route('/files/<int:corpID>', methods=['GET'])
def get_files_by_corp_id(corpID):
    matching_files = [{'doc_id': file.doc_id, 'file_name': file['file_name']} for file in fileTable.all() if file.get('company_code') == corpID]

    if not matching_files:
        return jsonify({"error": "No matched files"}), 404

    return jsonify({"files": matching_files}), 200

@QRapp.route('/companies', methods=['GET'])
def get_companies():
    start = int(request.args.get('start', 0))
    length = int(request.args.get('length', 10))
    search_value = request.args.get('search[value]', '').lower()

    companies = compTable.all()

    if search_value:
        companies = [comp for comp in companies if (
            search_value in comp['company'].lower() or
            search_value in comp.get('summary', '').lower() or
            search_value in comp.get('tags', '').lower()
        )]

    records_with_id = [{'doc_id': comp.doc_id, **comp} for comp in companies]
    records_with_id.reverse()
    paginated_companies = records_with_id[start:start + length]

    response = {
        "draw": int(request.args.get('draw', 1)),
        "recordsTotal": len(compTable.all()),
        "recordsFiltered": len(records_with_id),
        "data": paginated_companies
    }

    return jsonify(response)

@QRapp.route('/searchCompany', methods=['POST'])
def search_company():
    data = request.get_json()
    company_name = data.get('companyName').strip()

    if not company_name:
        return jsonify({'error': 'Company name cannot be empty!'}), 400
    
    results = None
    results = similar_search(company_name, threshold=80)

    return jsonify(results), 200

@QRapp.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

@QRapp.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    print(file.filename)

    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "File type not allowed"}), 400

    file_path = os.path.join(UPLOAD_FOLDER, file.filename)

    if os.path.exists(file_path):
        return jsonify({"error": "File already exists"}), 400

    try:
        file.save(file_path)
        text = extract_text_from_file(file_path)
        summary = ai_analysis(AI_PROMPT, text)

        fileTable.insert({
            'file_name': file.filename,
            'tags': '',
            'summary': summary,
            'location': file_path,
            'comments': '',
            'company': '',
            'company_code': ''
        })
        return jsonify({"success": "File uploaded", "filename": file.filename}), 201
    except Exception as e:
        logging.error(f"Error saving file: {e}")
        return jsonify({"error": "File could not be saved"}), 500

@QRapp.route('/aisummary', methods=['POST'])
def aisummary():
    data = request.json
    file_path = data.get('fileLocation')
    try:
        text = extract_text_from_file(file_path)
        prompt = "Summarize the content in the primary language used:"
        summary = ai_analysis(prompt, text)

        return jsonify({"success": "Successfully summarized", "summary": summary}), 201
    except Exception as e:
        logging.error(f"Error saving file: {e}")
        return jsonify({"error": "File could not be saved"}), 500

@QRapp.route('/delete', methods=['DELETE'])
def delete_file():
    doc_id = request.args.get('doc_id', type=int)
    
    if doc_id is None:
        return jsonify({"error": "No doc_id provided"}), 400

    file_entry = fileTable.get(doc_id=doc_id)

    if not file_entry:
        return jsonify({"error": "File not found"}), 404

    try:
        file_path = file_entry['location']
        fileTable.remove(doc_ids=[doc_id])
        os.remove(file_path)
        return jsonify({"success": "File deleted", "doc_id": doc_id}), 200
    except Exception as e:
        logging.error(f"Error deleting file: {e}")
        return jsonify({"error": "File could not be deleted"}), 500
    
@QRapp.route('/createcompany', methods=['POST'])
def create_company():
    data = request.json 
    companyName = data.get('companyName') 

    if not companyName:
        return jsonify({'error': 'Company name is required'}), 400

    new_id, companyData = create_new_company_data(companyName)

    if new_id is not None:
        return jsonify({'id': new_id, 'company': companyData['company']}), 201
    else:
        return jsonify({'error': 'Failed to create new company data'}), 500

def similar_search(company_name, threshold=80):
    results = compTable.all()    
    similar_results = [
        {'id': result.doc_id, 'company': result['company'], 'en_company': result['en_company']} 
        for result in results 
        if (fuzz.partial_ratio(company_name.lower(), result['company'].lower()) >= threshold or
            fuzz.partial_ratio(company_name.lower(), result['en_company'].lower()) >= threshold)
    ]    
    return similar_results

def allowed_file(filename):
    ALLOWED_EXTENSIONS = {'pdf', 'xls', 'xlsx', 'doc', 'docx', 'ppt', 'pptx'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    text = ''

    try:
        if ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + '\n'
        elif ext in ['.doc', '.docx']:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + '\n'
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
        elif ext in ['.xls', '.xlsx']:
            if ext == '.xls':
                workbook = xlrd.open_workbook(file_path)
                for sheet in workbook.sheets():
                    for row_idx in range(sheet.nrows):
                        row = sheet.row(row_idx)
                        text += ' '.join([str(cell.value) for cell in row]) + '\n'
            else:
                workbook = openpyxl.load_workbook(file_path, data_only=True)
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows(values_only=True):
                        text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
        elif ext == '.ppt':
            raise ValueError("PPT format is not supported. Please convert to PPTX format.")
    except Exception as e:
        logging.error(f"Error extracting text from file: {e}")
        raise  # 예외를 다시 발생시키면 호출자에게 알림

    return text

# def ai_analysis2(prompt, text):
#     try:
#         response = openai.chat.completions.create(
#             model="gpt-4o-mini",
#             messages=[
#                 {"role": "system", "content": "You are a helpful assistant that summarizes text."},
#                 {"role": "user", "content": f"{prompt}\n\n{text}"}
#             ],
#             max_tokens=4096
#         )
#         summary = response.choices[0].message.content.strip()
#         return summary
#     except Exception as e:
#         logging.error(f"Error during AI analysis: {e}")
#         return "Error during AI analysis."
    
def ai_analysis(prompt, text):

    url = 'https://api.openai.com/v1/chat/completions'
    
    headers = {
        'Content-Type': 'application/json',
        'Authorization': f'Bearer {QR_config['OPENAI_KEY']}'
    }

    data = {
        'model': "gpt-4o-mini",
        'messages': [
            {"role": "system", "content": "You are a helpful assistant that summarizes text."},
            {"role": "user", "content": f"{prompt}\n\n{text}"}
        ],
        'max_tokens': 4096,
        'temperature': 0.7
    }

    try:
        response = requests.post(url, headers=headers, json=data)
        response.raise_for_status()  # Raise an error on a bad response
        result = response.json()
        print('Response:', result)
        summary = result['choices'][0]['message']['content'].strip()
        return summary
    except requests.exceptions.HTTPError as err:
        print(f'Error: {err}')
        error_details = response.json()
        print('Error details:', error_details)
    
def create_new_company_data(companyName):
    try:
        companyData = {
            'company': companyName,
            'en_company': '',
            'industry': '',
            'country': '',
            'tags': '',
            'summary': '',
            'comments': ''
        }
        
        new_id = compTable.insert(companyData)
        return new_id, companyData
    
    except Exception as e:
        print(f"An error occurred: {e}")
        return None, None


if __name__ == '__main__':
    QRapp.run(host='0.0.0.0', port=QR_config['port'], debug=True)